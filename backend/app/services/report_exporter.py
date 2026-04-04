"""
Report Exporters
Converts structured report data into PDF, Markdown, and PPTX formats.
"""

import io
import os
import base64
import re
from typing import Dict, List, Any

# ═══════════════════════ MARKDOWN EXPORTER ═══════════════════════

def export_markdown(report: Dict[str, Any]) -> str:
    """Export report as a clean Markdown document."""
    lines = []

    for section in report.get('sections', []):
        s_type = section.get('type', 'text')
        title = section.get('title', '')

        if s_type == 'cover':
            lines.append(f"# {title}")
            if section.get('subtitle'):
                lines.append(f"\n**{section['subtitle']}**\n")
            if section.get('meta'):
                lines.append(f"*{section['meta']}*\n")
            lines.append("---\n")
            continue

        lines.append(f"\n## {title}\n")

        if s_type == 'text':
            lines.append(section.get('text', '') + "\n")

        elif s_type == 'stats_grid':
            # Metrics as a table
            metrics = section.get('metrics', [])
            if metrics:
                lines.append("| Metric | Value |")
                lines.append("|--------|-------|")
                for m in metrics:
                    lines.append(f"| {m['label']} | {m['value']} |")
                lines.append("")
            # Columns table
            cols_table = section.get('columns_table', [])
            if cols_table:
                lines.append("\n### Column Details\n")
                lines.append("| Column | Type | Nulls | Unique |")
                lines.append("|--------|------|-------|--------|")
                for c in cols_table:
                    lines.append(f"| {c['name']} | {c['dtype']} | {c['nulls']} | {c['unique']} |")
                lines.append("")

        elif s_type == 'quality_score':
            lines.append(f"**Overall Score: {section.get('score', 0)}/100 (Grade: {section.get('grade', 'N/A')})**\n")
            for b in section.get('breakdown', []):
                lines.append(f"- {b['label']}: {b['value']}%")
            lines.append("")

        elif s_type == 'analysis':
            if section.get('text'):
                lines.append(section['text'] + "\n")
            if section.get('chart'):
                lines.append(f"![{title}](data:image/png;base64,{section['chart']})\n")
            if section.get('table'):
                table = section['table']
                if table and isinstance(table[0], dict):
                    headers = list(table[0].keys())
                    lines.append("| " + " | ".join(headers) + " |")
                    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    for row in table:
                        lines.append("| " + " | ".join(str(row.get(h, '')) for h in headers) + " |")
                    lines.append("")
            if section.get('top_correlations'):
                lines.append("\n### Top Correlations\n")
                lines.append("| Feature 1 | Feature 2 | Correlation |")
                lines.append("|-----------|-----------|-------------|")
                for c in section['top_correlations'][:10]:
                    lines.append(f"| {c['col1']} | {c['col2']} | {c['value']:.4f} |")
                lines.append("")

        elif s_type == 'charts_grid':
            if section.get('stats_table'):
                lines.append("### Statistical Summary\n")
                st = section['stats_table']
                if st:
                    headers = list(st[0].keys())
                    lines.append("| " + " | ".join(headers) + " |")
                    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    for row in st:
                        lines.append("| " + " | ".join(str(row.get(h, '')) for h in headers) + " |")
                    lines.append("")
            for chart_item in section.get('charts', []):
                if chart_item.get('chart'):
                    lines.append(f"![{chart_item.get('title', '')}](data:image/png;base64,{chart_item['chart']})\n")

        elif s_type == 'model_table':
            if section.get('chart'):
                lines.append(f"![Model Comparison](data:image/png;base64,{section['chart']})\n")
            lines.append(f"**Best Model: {section.get('best_model', 'N/A')}**\n")
            models = section.get('models', [])
            if models:
                all_keys = set()
                for m in models:
                    all_keys.update(m.get('metrics', {}).keys())
                metric_keys = sorted(all_keys)[:6]
                headers = ['Model'] + metric_keys + ['Time (s)']
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for m in models:
                    prefix = "🏆 " if m.get('is_best') else ""
                    vals = [f"{prefix}{m['name']}"]
                    for k in metric_keys:
                        v = m.get('metrics', {}).get(k)
                        vals.append(f"{v:.4f}" if v is not None else "—")
                    vals.append(str(m.get('training_time', '')))
                    lines.append("| " + " | ".join(vals) + " |")
                lines.append("")

        elif s_type == 'insights':
            for insight in section.get('insights', []):
                lines.append(f"- {insight}")
            if section.get('ai_insights'):
                lines.append(f"\n### AI-Generated Insights\n\n{section['ai_insights']}\n")
            lines.append("")

        elif s_type == 'list':
            for item in section.get('items', []):
                lines.append(f"- {item}")
            lines.append("")

        elif s_type == 'ordered_list':
            for i, item in enumerate(section.get('items', []), 1):
                lines.append(f"{i}. {item}")
            lines.append("")

        elif s_type == 'log':
            if section.get('text'):
                lines.append(section['text'] + "\n")
            for step in section.get('steps', []):
                cat = step.get('category', '').upper()
                lines.append(f"- **[{cat}]** {step['step']}")
            lines.append("")

        elif s_type == 'kpi_cards':
            for card in section.get('cards', []):
                lines.append(f"- **{card['label']}**: {card['value']}")
            lines.append("")

    return "\n".join(lines)


# ═══════════════════════ PDF EXPORTER ═══════════════════════

def export_pdf(report: Dict[str, Any]) -> bytes:
    """Export report as a multi-page PDF using fpdf2."""
    from fpdf import FPDF

    def clean_text(text):
        """Remove markdown bold/italic markers for PDF."""
        text = str(text) if text is not None else ""
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        # Replace special chars
        text = text.replace('↔', '<->').replace('—', '-')
        return text.encode('latin-1', 'replace').decode('latin-1')

    class ReportPDF(FPDF):
        def __init__(self):
            super().__init__()
            self.set_auto_page_break(auto=True, margin=20)
            self.report_title = ''

        def header(self):
            if self.page_no() > 1:
                self.set_font('Helvetica', 'I', 8)
                self.set_text_color(140, 140, 140)
                self.cell(0, 8, clean_text(self.report_title), align='L')
                self.cell(0, 8, 'ML Yantra', align='R')
                self.ln(12)
                self.set_draw_color(220, 220, 220)
                self.line(10, self.get_y(), 200, self.get_y())
                self.ln(4)

        def footer(self):
            self.set_y(-15)
            self.set_font('Helvetica', 'I', 8)
            self.set_text_color(160, 160, 160)
            self.cell(0, 10, f'Page {self.page_no()}/{{nb}}', align='C')

    pdf = ReportPDF()
    pdf.alias_nb_pages()
    pdf.report_title = report.get('title', 'Report')

    def add_image_from_b64(b64_str, w=170):
        if not b64_str:
            return
        try:
            img_data = base64.b64decode(b64_str)
            img_stream = io.BytesIO(img_data)
            pdf.image(img_stream, x=20, w=w)
            pdf.ln(5)
        except Exception:
            pass

    for section in report.get('sections', []):
        s_type = section.get('type', 'text')
        title = section.get('title', '')

        if s_type == 'cover':
            if pdf.page_no() == 0 or pdf.get_y() > 35:
                pdf.add_page()
            pdf.ln(50)
            pdf.set_font('Helvetica', 'B', 28)
            pdf.set_text_color(171, 53, 5)
            pdf.multi_cell(0, 14, clean_text(title), align='C')
            pdf.ln(8)
            if section.get('subtitle'):
                pdf.set_font('Helvetica', '', 14)
                pdf.set_text_color(80, 80, 80)
                pdf.multi_cell(0, 8, clean_text(section['subtitle']), align='C')
            pdf.ln(15)
            if section.get('meta'):
                pdf.set_font('Helvetica', 'I', 10)
                pdf.set_text_color(140, 140, 140)
                pdf.multi_cell(0, 7, clean_text(section['meta']), align='C')
            continue

        if pdf.page_no() == 0 or pdf.get_y() > 35:
            pdf.add_page()
        # Section title
        pdf.set_font('Helvetica', 'B', 16)
        pdf.set_text_color(171, 53, 5)
        pdf.cell(0, 10, clean_text(title))
        pdf.ln(12)
        pdf.set_text_color(40, 40, 40)

        if s_type in ('text', 'analysis'):
            if section.get('text'):
                pdf.set_font('Helvetica', '', 10)
                pdf.multi_cell(0, 6, clean_text(section['text']))
                pdf.ln(5)
            if section.get('chart'):
                add_image_from_b64(section['chart'])
            # Table
            table = section.get('table', [])
            if table and isinstance(table, list) and isinstance(table[0], dict):
                headers = list(table[0].keys())
                col_w = 170 / len(headers)
                pdf.set_font('Helvetica', 'B', 8)
                pdf.set_fill_color(245, 245, 245)
                for h in headers:
                    pdf.cell(col_w, 7, clean_text(str(h)), border=1, fill=True, align='C')
                pdf.ln()
                pdf.set_font('Helvetica', '', 8)
                for row in table[:20]:
                    for h in headers:
                        pdf.cell(col_w, 6, clean_text(str(row.get(h, ''))), border=1, align='C')
                    pdf.ln()
                pdf.ln(5)
            # Top correlations
            top_corrs = section.get('top_correlations', [])
            if top_corrs:
                pdf.set_font('Helvetica', 'B', 11)
                pdf.cell(0, 8, 'Top Correlations')
                pdf.ln(8)
                pdf.set_font('Helvetica', 'B', 8)
                pdf.set_fill_color(245, 245, 245)
                for h in ['Feature 1', 'Feature 2', 'Correlation']:
                    pdf.cell(56, 7, h, border=1, fill=True, align='C')
                pdf.ln()
                pdf.set_font('Helvetica', '', 8)
                for c in top_corrs[:10]:
                    pdf.cell(56, 6, clean_text(str(c['col1'])), border=1, align='C')
                    pdf.cell(56, 6, clean_text(str(c['col2'])), border=1, align='C')
                    pdf.cell(56, 6, f"{c['value']:.4f}", border=1, align='C')
                    pdf.ln()

        elif s_type == 'stats_grid':
            metrics = section.get('metrics', [])
            if metrics:
                pdf.set_font('Helvetica', 'B', 8)
                pdf.set_fill_color(245, 245, 245)
                pdf.cell(85, 7, 'Metric', border=1, fill=True, align='C')
                pdf.cell(85, 7, 'Value', border=1, fill=True, align='C')
                pdf.ln()
                pdf.set_font('Helvetica', '', 9)
                for m in metrics:
                    pdf.cell(85, 6, clean_text(m['label']), border=1, align='L')
                    pdf.cell(85, 6, clean_text(str(m['value'])), border=1, align='C')
                    pdf.ln()
            pdf.ln(5)
            cols_table = section.get('columns_table', [])
            if cols_table:
                pdf.set_font('Helvetica', 'B', 11)
                pdf.cell(0, 8, 'Column Details')
                pdf.ln(8)
                pdf.set_font('Helvetica', 'B', 7)
                pdf.set_fill_color(245, 245, 245)
                for h in ['Column', 'Type', 'Nulls', 'Unique']:
                    pdf.cell(42, 6, h, border=1, fill=True, align='C')
                pdf.ln()
                pdf.set_font('Helvetica', '', 7)
                for c in cols_table[:25]:
                    pdf.cell(42, 5, clean_text(str(c['name'])[:20]), border=1)
                    pdf.cell(42, 5, clean_text(str(c['dtype'])), border=1, align='C')
                    pdf.cell(42, 5, str(c['nulls']), border=1, align='C')
                    pdf.cell(42, 5, str(c['unique']), border=1, align='C')
                    pdf.ln()

        elif s_type == 'quality_score':
            pdf.set_font('Helvetica', 'B', 24)
            pdf.set_text_color(171, 53, 5)
            pdf.cell(0, 15, clean_text(f"Score: {section.get('score', 0)}/100 ({section.get('grade', '?')})"), align='C')
            pdf.ln(18)
            pdf.set_text_color(40, 40, 40)
            pdf.set_font('Helvetica', '', 10)
            for b in section.get('breakdown', []):
                pdf.cell(0, 7, clean_text(f"{b['label']}: {b['value']}%"))
                pdf.ln()

        elif s_type == 'charts_grid':
            # Stats table
            st = section.get('stats_table', [])
            if st:
                pdf.set_font('Helvetica', 'B', 7)
                pdf.set_fill_color(245, 245, 245)
                headers = list(st[0].keys())
                col_w = 170 / len(headers)
                for h in headers:
                    pdf.cell(col_w, 6, clean_text(h), border=1, fill=True, align='C')
                pdf.ln()
                pdf.set_font('Helvetica', '', 7)
                for row in st:
                    for h in headers:
                        pdf.cell(col_w, 5, clean_text(str(row.get(h, ''))), border=1, align='C')
                    pdf.ln()
                pdf.ln(5)
            for chart_item in section.get('charts', []):
                if chart_item.get('chart'):
                    add_image_from_b64(chart_item['chart'], w=150)

        elif s_type == 'model_table':
            if section.get('chart'):
                add_image_from_b64(section['chart'])
            pdf.set_font('Helvetica', 'B', 11)
            pdf.cell(0, 8, f"Best Model: {clean_text(section.get('best_model', 'N/A'))}")
            pdf.ln(10)
            models = section.get('models', [])
            if models:
                all_keys = set()
                for m in models:
                    all_keys.update(m.get('metrics', {}).keys())
                metric_keys = sorted(all_keys)[:5]
                headers = ['Model'] + metric_keys
                col_w = 170 / len(headers)
                pdf.set_font('Helvetica', 'B', 7)
                pdf.set_fill_color(245, 245, 245)
                for h in headers:
                    pdf.cell(col_w, 6, clean_text(h)[:12], border=1, fill=True, align='C')
                pdf.ln()
                pdf.set_font('Helvetica', '', 7)
                for m in models:
                    name = ('* ' if m.get('is_best') else '') + m['name']
                    pdf.cell(col_w, 5, clean_text(name)[:15], border=1)
                    for k in metric_keys:
                        v = m.get('metrics', {}).get(k)
                        pdf.cell(col_w, 5, f"{v:.4f}" if v is not None else "-", border=1, align='C')
                    pdf.ln()

        elif s_type == 'insights':
            pdf.set_font('Helvetica', '', 10)
            for insight in section.get('insights', []):
                pdf.multi_cell(0, 6, f"  * {clean_text(insight)}")
                pdf.ln(2)
            if section.get('ai_insights'):
                pdf.ln(4)
                pdf.set_font('Helvetica', 'B', 11)
                pdf.cell(0, 8, 'AI-Generated Insights')
                pdf.ln(8)
                pdf.set_font('Helvetica', '', 10)
                pdf.multi_cell(0, 6, clean_text(section['ai_insights']))

        elif s_type in ('list', 'ordered_list'):
            pdf.set_font('Helvetica', '', 10)
            for i, item in enumerate(section.get('items', []), 1):
                prefix = f"{i}. " if s_type == 'ordered_list' else "  * "
                pdf.multi_cell(0, 6, f"{prefix}{clean_text(item)}")
                pdf.ln(2)

        elif s_type == 'log':
            if section.get('text'):
                pdf.set_font('Helvetica', '', 10)
                pdf.multi_cell(0, 6, clean_text(section['text']))
                pdf.ln(4)
            pdf.set_font('Helvetica', '', 8)
            for step in section.get('steps', []):
                cat = step.get('category', '').upper()
                pdf.multi_cell(0, 5, f"  [{cat}] {clean_text(step['step'])}")
                pdf.ln(1)

        elif s_type == 'kpi_cards':
            pdf.set_font('Helvetica', '', 11)
            for card in section.get('cards', []):
                pdf.cell(0, 8, clean_text(f"  {card['label']}: {card['value']}"))
                pdf.ln()

    return pdf.output()


# ═══════════════════════ PPTX EXPORTER ═══════════════════════

def export_pptx(report: Dict[str, Any]) -> bytes:
    """Export report as a PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BRAND = RGBColor(171, 53, 5)
    WHITE = RGBColor(255, 255, 255)
    DARK = RGBColor(30, 30, 30)
    GRAY = RGBColor(120, 120, 120)

    def add_text_box(slide, left, top, width, height, text, font_size=14,
                     bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return tf

    def add_image_from_b64(slide, b64_str, left, top, width):
        if not b64_str:
            return
        try:
            img_data = base64.b64decode(b64_str)
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(img_stream, Inches(left), Inches(top), Inches(width))
        except Exception:
            pass

    def clean_md(text):
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        return text

    for section in report.get('sections', []):
        s_type = section.get('type', 'text')
        title = section.get('title', '')

        if s_type == 'cover':
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            # Background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(26, 28, 26)
            # Title
            add_text_box(slide, 1, 2, 11, 1.5, clean_md(title), font_size=40, bold=True, color=BRAND, alignment=PP_ALIGN.CENTER)
            if section.get('subtitle'):
                add_text_box(slide, 1, 3.8, 11, 0.8, clean_md(section['subtitle']), font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
            if section.get('meta'):
                add_text_box(slide, 1, 5, 11, 0.6, clean_md(section['meta']), font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
            continue

        # Content slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Title bar
        add_text_box(slide, 0.5, 0.3, 12, 0.6, clean_md(title), font_size=24, bold=True, color=BRAND)

        if s_type in ('text', 'analysis'):
            y_pos = 1.2
            if section.get('text'):
                add_text_box(slide, 0.7, y_pos, 11.5, 1.5, clean_md(section['text']), font_size=13, color=DARK)
                y_pos += 1.8
            if section.get('chart'):
                add_image_from_b64(slide, section['chart'], 1.5, y_pos, 9)

        elif s_type == 'stats_grid':
            metrics = section.get('metrics', [])
            for i, m in enumerate(metrics[:8]):
                col = i % 4
                row = i // 4
                x = 0.8 + col * 3.1
                y = 1.3 + row * 1.5
                add_text_box(slide, x, y, 2.8, 0.4, m['value'], font_size=28, bold=True, color=BRAND, alignment=PP_ALIGN.CENTER)
                add_text_box(slide, x, y + 0.5, 2.8, 0.3, m['label'], font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

        elif s_type == 'quality_score':
            add_text_box(slide, 2, 2, 9, 1, f"{section.get('score', 0)}/100", font_size=60, bold=True, color=BRAND, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 2, 3.5, 9, 0.6, f"Grade: {section.get('grade', '?')}", font_size=24, color=DARK, alignment=PP_ALIGN.CENTER)

        elif s_type == 'charts_grid':
            charts = section.get('charts', [])
            for i, chart_item in enumerate(charts[:4]):
                col = i % 2
                row = i // 2
                x = 0.5 + col * 6.2
                y = 1.2 + row * 3
                if chart_item.get('chart'):
                    add_image_from_b64(slide, chart_item['chart'], x, y, 5.8)

        elif s_type == 'model_table':
            if section.get('chart'):
                add_image_from_b64(slide, section['chart'], 1.5, 1.2, 9)
            add_text_box(slide, 0.7, 5.5, 11, 0.5, f"Best Model: {clean_md(section.get('best_model', 'N/A'))}", font_size=16, bold=True, color=BRAND)

        elif s_type == 'insights':
            y = 1.3
            for insight in section.get('insights', []):
                add_text_box(slide, 0.7, y, 11.5, 0.5, f"• {clean_md(insight)}", font_size=13, color=DARK)
                y += 0.55

        elif s_type in ('list', 'ordered_list'):
            y = 1.3
            for i, item in enumerate(section.get('items', []), 1):
                prefix = f"{i}. " if s_type == 'ordered_list' else "• "
                add_text_box(slide, 0.7, y, 11.5, 0.5, f"{prefix}{clean_md(item)}", font_size=13, color=DARK)
                y += 0.55

        elif s_type == 'log':
            y = 1.3
            if section.get('text'):
                add_text_box(slide, 0.7, y, 11.5, 0.5, clean_md(section['text']), font_size=13, color=DARK)
                y += 0.7
            for step in section.get('steps', [])[:12]:
                cat = step.get('category', '').upper()
                add_text_box(slide, 0.7, y, 11.5, 0.4, f"[{cat}] {clean_md(step['step'])}", font_size=10, color=GRAY)
                y += 0.4

        elif s_type == 'kpi_cards':
            cards = section.get('cards', [])
            for i, card in enumerate(cards[:5]):
                x = 0.5 + i * 2.5
                add_text_box(slide, x, 2, 2.2, 0.6, str(card['value']), font_size=30, bold=True, color=BRAND, alignment=PP_ALIGN.CENTER)
                add_text_box(slide, x, 2.8, 2.2, 0.4, card['label'], font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
